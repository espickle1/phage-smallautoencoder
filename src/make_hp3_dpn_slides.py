#!/usr/bin/env python3
"""Build the HP3 gp8 DPN-duplication PowerPoint deck (3 slides).

Slide 1: analysis/figures/cluster_structures_hp3/gp8_2mer_evolved_roi_clusters.png --
    evolved gp8 homodimer, SAE-cluster colored, ROI 287-293 marked. Paired with a
    table summarizing the SAE feature-rank findings for that ROI (WT vs evolved).
Slide 2: analysis/figures/cluster_structures_hp3/gp8_evolved_in_wedge_context.png --
    same gp8, placed in its real structural neighborhood (T4 baseplate wedge, 9MKB).
Slide 3: analysis/figures/cluster_structures_hp3/t4_tail_overlay_montage_3x3.png --
    3x3 montage (T4/HP3/HP3e x axial/transverse/cross-section) of gp8 overlaid on
    the real T4 tail.

Each slide pairs its figure with short on-slide bullets; extended detail (rank 4-5
nuance, background context) goes into that slide's speaker notes.

Usage:
    python src/make_hp3_dpn_slides.py
"""
from __future__ import annotations

import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIGURES = ROOT / "analysis" / "figures" / "cluster_structures_hp3"
OUT = ROOT / "manuscript" / "slides" / "hp3_gp8_dpn_duplication.pptx"

INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x89, 0x87, 0x81)
GRID = RGBColor(0xE1, 0xE0, 0xD9)
HEADER_FILL = RGBColor(0xE1, 0xE0, 0xD9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.4)
FIG_TOP = Inches(1.25)
FIG_BOTTOM_MAX = Inches(7.1)

TABLE_HEADERS = [
    "Region", "Residues (WT)", "Residues (Evolved)",
    "Rank 1 (Disorder)", "Rank 2 (Struct. motif)", "Rank 3",
]
TABLE_ROWS = [
    ["First DPN copy (invariant)", "D287-P288-N289", "D287-P288-N289",
     "15659 (~0.60-0.65)", "15513 (~0.40-0.44)", "7819 -- Interaction site"],
    ["V290 / relocated 2nd copy", "V290", "D290-P291-N292-V293",
     "15659 (~0.60-0.65)", "15513 (~0.40-0.44)", "3038 -- Domain"],
]

SLIDE_1 = dict(
    title="Evolved gp8 ROI (287-293): SAE feature summary",
    panel=FIGURES / "gp8_2mer_evolved_roi_clusters.png",
    fig_height=Inches(5.6),
    body_size=Pt(13),
    intro=[
        "Ranks 1-2 are identical across all of 287-293, in both variants: dominant "
        "“whole-chain activation / central peaks” (Disorder, feature 15659, "
        "~0.60–0.65) plus a weaker “catalytic/interface βαβ module” "
        "(Structural motif, feature 15513, ~0.40–0.44) — this region reads "
        "primarily as a flexible surface loop.",
        "Rank 3 is where the two variants diverge — but at the *same sequence "
        "position* in both: the first DPN copy (D287-P288-N289) always reads “Phage "
        "and viral structural proteins” (Interaction site, feature 7819); V290 (WT) "
        "and the inserted second copy D290-P291-N292-V293 (evolved) both read "
        "“Extended structured domain surfaces” (Domain, feature 3038) instead.",
        "Conclusion: the duplication *relocates* the interaction site rather than "
        "duplicating it — ESM-C/SAE only ever calls one DPN copy "
        "interaction-relevant, and it's always the first one.",
    ],
    caveat=(
        "ESM-C is sequence-only — these are learned sequence-context labels, not "
        "measured 3D contacts (see Slides 2–3)."
    ),
    legend=(
        "Speaker notes. Evolved 290-293 rank 4 is feature 7819 (Interaction site) on all "
        "four residues (~0.34-0.37) -- a secondary echo of the interaction signal, not a "
        "second strong site. Rank 5 varies per-residue: D290 -> 7070 \"Phage tail "
        "attachment segments\" (Disorder, 0.295); P291 -> 3509 \"C-terminal helix-to-IDR "
        "transition\" (Disorder, 0.303); N292 -> 12804 \"Cysteine-rich LU/CRD domains\" "
        "(Domain, 0.303); V293 -> 3509 (Disorder, 0.316). WT V290 rank 4/5 = 7819 (0.359) "
        "/ 7070 (0.309)."
    ),
)

SLIDE_2 = dict(
    title="gp8 (evolved) in its real T4 wedge context",
    panel=FIGURES / "gp8_evolved_in_wedge_context.png",
    fig_height=Inches(5.6),
    body_size=Pt(14),
    intro=[
        "Same evolved gp8, SAE-cluster colored, now shown inside its authentic "
        "structural neighborhood: the cryo-EM T4 baseplate wedge (PDB 9MKB), not an "
        "isolated monomer.",
        "Confirms gp8 is a genuine T4 gp8 baseplate-wedge homolog — CA-centroid "
        "overlay onto 9MKB chain ZE differs by only 0.30 Å.",
        "Wedge stoichiometry recovered exactly as expected from T4 biology: 2x gp8 + "
        "gp7 + gp6 (pale blue / pink / tan, de-emphasized as background context).",
        "The ROI (287-293, black spheres) sits right at the boundary facing the "
        "second gp8 copy — i.e., near a real inter-subunit interface, not buried "
        "mid-fold.",
        "Open question this figure sets up: does the newly-inserted second DPN copy "
        "make its own contact with a neighboring chain, or is it simply repositioned "
        "in place? Requires a direct contact measurement from the aligned-overlay "
        "structure — next planned step, not done yet.",
    ],
    caveat=None,
    legend=(
        "Speaker notes. 9MKB is the cryo-EM structure of the bacteriophage T4 "
        "portal-neck-tail complex (541 chains). A CA-CA contact search (<10 Å) from "
        "chain ZE (which overlays our AF3 gp8 model almost exactly) against all 541 "
        "other chains finds exactly 6 true neighbors: ZD (2nd gp8 copy), YC/Yc (gp7), "
        "eB/eU/f9 (gp6) -- matching the known T4 wedge stoichiometry (2x gp8 + gp7 + gp6 "
        "per wedge). All other 534 chains (tail tube/sheath, capsid, portal, fibers) are "
        "irrelevant to gp8's immediate context and were dropped, not just hidden."
    ),
)

SLIDE_3 = dict(
    title="T4 tail overlay: does the duplication move the contact point?",
    panel=FIGURES / "t4_tail_overlay_montage_3x3.png",
    fig_height=Inches(5.1),
    body_size=Pt(13),
    intro=[
        "3x3 comparison: columns = T4 gp8 (reference), HP3 (WT), HP3e (evolved); "
        "rows = axial, transverse, and cross-section views of gp8 overlaid on the "
        "real T4 tail.",
        "Across all three viewing angles, the DPN duplication visibly shifts the "
        "first-contact loop's position/orientation relative to WT — independent, "
        "structural confirmation of the Slide 1 SAE finding: the interaction "
        "footprint doesn't get bigger, it moves.",
        "Why duplicate the whole DPN triplet, and not a partial insertion (DP, P, or "
        "PN alone)? DPN is a proline-anchored turn unit — a partial insertion "
        "would break the local turn geometry, while duplicating the full triplet "
        "preserves the fold and just shifts its downstream position.",
        "Working conclusion: this duplication event *repositions the binding "
        "location*, consistent with the tropism shift observed in HP3e.",
        "Still open: whether the inserted second copy independently contacts the "
        "tail (rather than just riding along structurally) isn't resolved by either "
        "the SAE features (sequence-only) or this overlay — flagged for the "
        "follow-up contact analysis.",
    ],
    caveat=None,
    legend=(
        "Speaker notes. Panels built from AF3 models of HP3 (WT) and HP3e (evolved) gp8, "
        "structurally aligned onto the real T4 tail (9MKB) coordinate frame, then "
        "rendered from three fixed camera angles per variant so the same viewpoint is "
        "directly comparable across T4/HP3/HP3e."
    ),
)


def _add_title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.25), Inches(12.5), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = INK

    line = slide.shapes.add_shape(1, MARGIN, Inches(1.05), Inches(12.5), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = GRID
    line.line.fill.background()


def _add_run_with_italics(p, para_text: str, size: Pt, color: RGBColor) -> None:
    # *word* -> italic run, so emphasis renders instead of literal asterisks
    for j, segment in enumerate(re.split(r"\*(.+?)\*", para_text)):
        if not segment:
            continue
        run = p.add_run()
        run.text = segment
        run.font.size = size
        run.font.color.rgb = color
        run.font.italic = bool(j % 2)


def _add_figure(slide, spec: dict) -> Inches:
    im = Image.open(spec["panel"])
    fig_h = spec["fig_height"]
    fig_w = Inches(fig_h.inches * im.width / im.height)
    fig_h = min(fig_h, FIG_BOTTOM_MAX - FIG_TOP)
    slide.shapes.add_picture(str(spec["panel"]), MARGIN, FIG_TOP, height=fig_h, width=fig_w)
    return fig_w


def _add_caption(slide, spec: dict, cap_left: Inches, cap_top: Inches, cap_width: Inches,
                  cap_height: Inches) -> None:
    cap = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
    ctf = cap.text_frame
    ctf.word_wrap = True
    for i, para_text in enumerate(spec["intro"]):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        _add_run_with_italics(p, para_text, spec["body_size"], INK)
    if spec.get("caveat"):
        p = ctf.add_paragraph()
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = spec["caveat"]
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = MUTED


def _add_feature_table(slide, left: Inches, top: Inches, width: Inches, height: Inches) -> None:
    n_rows = 1 + len(TABLE_ROWS)
    n_cols = len(TABLE_HEADERS)
    gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gfx.table

    # narrower Region/Residues columns, wider rank columns
    col_fracs = [0.16, 0.15, 0.19, 0.17, 0.17, 0.16]
    for idx, frac in enumerate(col_fracs):
        table.columns[idx].width = Emu(int(width * frac))

    for c, text in enumerate(TABLE_HEADERS):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = INK

    for r, row in enumerate(TABLE_ROWS, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = INK


def _build_slide(prs: Presentation, spec: dict, with_table: bool = False) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _add_title(slide, spec["title"])
    fig_w = _add_figure(slide, spec)

    cap_left = MARGIN + fig_w + Inches(0.35)
    cap_width = SLIDE_W - cap_left - MARGIN

    if with_table:
        cap_height = Inches(3.0)
        _add_caption(slide, spec, cap_left, FIG_TOP, cap_width, cap_height)
        table_top = FIG_TOP + cap_height + Inches(0.15)
        table_height = Inches(1.6)
        _add_feature_table(slide, cap_left, table_top, cap_width, table_height)
    else:
        cap_height = FIG_BOTTOM_MAX - FIG_TOP
        _add_caption(slide, spec, cap_left, FIG_TOP, cap_width, cap_height)

    slide.notes_slide.notes_text_frame.text = spec["legend"]


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_slide(prs, SLIDE_1, with_table=True)
    _build_slide(prs, SLIDE_2)
    _build_slide(prs, SLIDE_3)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
