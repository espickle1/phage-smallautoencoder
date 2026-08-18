#!/usr/bin/env python3
"""Build the HP3 tail-collar-fiber (TCF) + baseplate-gp8 report deck.

Presentation style follows the Maresso-group / Baylor College of Medicine house
template (title slide with gray banner + Baylor logo; content slides with dual
navy header bars, an italic centered subtitle, a gray citation footer, and page
numbers), reproduced from manuscript/slides/Protein_Structure_Group_Meeting.

The deck order:

  Slide 1   TITLE (gray banner, presenter block, Baylor logo)
  Slide 2   Opening the black box — SAEs as a validated interpretability method
  Slide 3   Same method, our proteins (ESMC); labels come in post-hoc
  Slide 4   HP3 -> HP3.1 overview (+ annotated T4 schematic)
  Slide 5   TCF trimer: three substitutions, all at inter-chain contacts
  Slide 6   TCF SAE features: substitutions leave the feature profile unchanged
  Slide 7   TCF ROI down the 3-fold axis (head-domain surface derivative)
  Slide 8   gp8: a DPN tandem-repeat insertion relocates the interaction site
  Slide 9   gp8 in its real T4 baseplate-wedge context (9MKB)
  Slide 10  gp8 on the T4 tail: does the duplication move the contact point?
  Slide 11  Synthesis: substitution-at-a-conserved-site vs insertion-and-relocation

The two opener slides are generated natively in the same house style as the
rest of the deck (dual navy header bars, footer, etc.) rather than imported
from an external file, so the whole deck builds from this one script.

Usage:
    python src/make_hp3_tcf_gp8_report.py
"""
from __future__ import annotations

import csv
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "analysis" / "figures"
CS = FIG / "cluster_structures_hp3"
RS = FIG / "roi_structures_hp3"
ASSETS = ROOT / "manuscript" / "slides" / "assets"
BAYLOR_LOGO = ASSETS / "baylor_logo.png"
T4_OVERVIEW = FIG / "t4_overview_annotated.png"
TCF_HEAD_AXIS = CS / "tcf_3mer_head_axis_derivative.png"
TCF_HEAD_VIEWS = [
    (CS / "tcf_head_hp3e_full_on.png", "Down the 3-fold axis"),
    (CS / "tcf_head_hp3e_tilted.png", "Tilted"),
    (CS / "tcf_head_hp3e_side.png", "Side"),
]
TCF_TABLE_CSV = ROOT / "analysis" / "features" / "roi_feature_summary_hp3_tcf.csv"

# --- aspect ratio: "16:9" (default, 13.333x7.5) or "4:3" (10x7.5) ---
# Override from the environment: ASPECT=4:3 python src/make_hp3_tcf_gp8_report.py
import os as _os
ASPECT = _os.environ.get("ASPECT", "16:9")
if ASPECT not in ("16:9", "4:3"):
    raise SystemExit(f"ASPECT must be '16:9' or '4:3', got {ASPECT!r}")
_SUFFIX = "" if ASPECT == "16:9" else "_4x3"
OUT = ROOT / "manuscript" / "slides" / f"hp3_tcf_gp8_report{_SUFFIX}.pptx"

# --- presenter block on the title slide (edit here) ---
DECK_TITLE = "Reading Amino-Acid Function with Sparse Autoencoders"
DECK_SUBTITLE = "HP3 phage tail collar fiber and baseplate gp8"
PRESENTER = "James Chang, Ph.D."
GROUP = "Maresso Group"
DECK_DATE = "July 24, 2026"

# --- palette ---
INK = RGBColor(0x0B, 0x0B, 0x0B)
MUTED = RGBColor(0x89, 0x87, 0x81)
GRID = RGBColor(0xE1, 0xE0, 0xD9)
HEADER_FILL = RGBColor(0xE1, 0xE0, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# house-style institutional colors (sampled from the reference deck)
NAVY = RGBColor(0x00, 0x20, 0x60)         # header bars
BANNER_GRAY = RGBColor(0x59, 0x59, 0x59)  # title-slide banner
BAYLOR_NAVY = RGBColor(0x1A, 0x28, 0x61)  # logo navy
TCF_ACCENT = RGBColor(0x2A, 0x62, 0x8A)   # deep teal-blue for TCF section
GP8_ACCENT = RGBColor(0x8A, 0x3A, 0x2A)   # warm brick for gp8 section
DARK_BG = RGBColor(0x0B, 0x14, 0x2E)      # opener slide 1: "black box" ground
DIM_CELL = RGBColor(0x2A, 0x35, 0x54)     # opener: unlit feature-grid cell
LIT_CELL = RGBColor(0xE0, 0xA5, 0x3C)     # opener: lit ("activated") feature-grid cell
LIGHT_CELL = RGBColor(0xD8, 0xDC, 0xE6)   # opener slide 2: unlit cell on light ground

# The layout was designed on a 13.333in-wide (16:9) canvas. For 4:3 we keep the
# height fixed at 7.5in and rescale every HORIZONTAL coordinate by the width
# ratio, so the vertical rhythm (header bars, figure band, footer) is untouched.
_DESIGN_W = 13.333
SLIDE_W = Inches(13.333 if ASPECT == "16:9" else 10.0)
SLIDE_H = Inches(7.5)
WSCALE = SLIDE_W.inches / _DESIGN_W   # 1.0 for 16:9, 0.75 for 4:3


def HX(inches: float) -> Inches:
    """Scale a design-space horizontal measurement to the current aspect."""
    return Inches(inches * WSCALE)


MARGIN = HX(0.4)
# content-area band: below the navy header bars + subtitle, above the footer
BAR_H = Inches(0.95)
SUBTITLE_TOP = Inches(1.02)
FIG_TOP = Inches(1.75)
FIG_BOTTOM_MAX = Inches(6.85)
FOOTER_Y = Inches(7.02)


# ---------------------------------------------------------------------------
# house-style primitives (title banner, dual navy header bars, footer)
# ---------------------------------------------------------------------------
def _small_caps(run) -> None:
    """Render a run in small caps (Title Case -> large initial + small caps)."""
    run._r.get_or_add_rPr().set("cap", "small")


def _rect(slide, left, top, width, height, fill: RGBColor):
    shp = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_header_bars(slide, topic: str, section: str, accent: RGBColor) -> None:
    """Two navy bars: left = slide topic, right = section tag (house style)."""
    left_w = HX(6.9)
    right_x = HX(9.15)
    right_w = SLIDE_W - right_x
    # left bar (topic)
    bar_l = _rect(slide, Inches(0), Inches(0), left_w, BAR_H, NAVY)
    tf = bar_l.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = topic
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = WHITE
    _small_caps(r)
    # right bar (section)
    bar_r = _rect(slide, right_x, Inches(0), right_w, BAR_H, NAVY)
    tf = bar_r.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.25)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = section
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = WHITE
    _small_caps(r)


def _add_subtitle(slide, text: str, color: RGBColor = INK) -> None:
    """Italic centered subtitle just below the header bars."""
    tb = slide.shapes.add_textbox(MARGIN, SUBTITLE_TOP, SLIDE_W - 2 * MARGIN, Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = color


def _add_footer(slide, citation: str | None, page: int | None) -> None:
    """Gray citation bottom-left, page number bottom-right (house style)."""
    if citation:
        tb = slide.shapes.add_textbox(MARGIN, FOOTER_Y, SLIDE_W - HX(2.8), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        _add_run_with_italics(p, citation, Pt(11), MUTED)
    if page is not None:
        tb = slide.shapes.add_textbox(SLIDE_W - HX(0.93), FOOTER_Y, HX(0.7), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(page)
        r.font.size = Pt(12)
        r.font.color.rgb = MUTED


def slide_title(prs) -> None:
    """House-style title slide: gray banner, presenter block, Baylor logo."""
    slide = _blank_slide(prs)
    # gray banner across the top ~2/3 of the slide
    banner_h = Inches(5.03)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, banner_h, BANNER_GRAY)
    # title (white, small caps, in the banner)
    tb = slide.shapes.add_textbox(HX(1.0), Inches(1.15), SLIDE_W - HX(2.0), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = DECK_TITLE
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = WHITE
    _small_caps(r)
    # italic content subtitle line under the title
    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    r2 = p2.add_run()
    r2.text = DECK_SUBTITLE
    r2.font.size = Pt(20)
    r2.font.italic = True
    r2.font.color.rgb = WHITE
    # presenter block (name bold, group, date)
    pb = slide.shapes.add_textbox(HX(1.0), Inches(3.55), SLIDE_W - HX(2.0), Inches(1.3))
    ptf = pb.text_frame
    ptf.word_wrap = True
    for i, (txt, bold) in enumerate([(PRESENTER, True), (GROUP, False), (DECK_DATE, False)]):
        p = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(22) if i == 0 else Pt(20)
        r.font.bold = bold
        r.font.color.rgb = WHITE
    # Baylor logo, bottom-left on the white band
    if BAYLOR_LOGO.exists():
        im = Image.open(BAYLOR_LOGO)
        logo_h = Inches(1.7)
        logo_w = Inches(logo_h.inches * im.width / im.height)
        slide.shapes.add_picture(str(BAYLOR_LOGO), HX(0.55),
                                 SLIDE_H - logo_h - Inches(0.25),
                                 width=logo_w, height=logo_h)


def _add_run_with_italics(p, para_text: str, size, color: RGBColor) -> None:
    # *word* -> italic run, so emphasis renders instead of literal asterisks
    for j, segment in enumerate(re.split(r"\*(.+?)\*", para_text)):
        if not segment:
            continue
        run = p.add_run()
        run.text = segment
        run.font.size = size
        run.font.color.rgb = color
        run.font.italic = bool(j % 2)


def _fit_figure(panel: Path, target_h: Inches, max_w: Inches):
    im = Image.open(panel)
    h = min(target_h, FIG_BOTTOM_MAX - FIG_TOP)
    # width caps are given in 16:9 design inches; narrow them for 4:3 so the
    # figure still leaves room for the caption/table column beside it.
    max_w = Inches(max_w.inches * WSCALE)
    w = Inches(h.inches * im.width / im.height)
    if w > max_w:
        w = max_w
        h = Inches(w.inches * im.height / im.width)
    return w, h


def _add_figure(slide, panel: Path, target_h: Inches, left: Inches = MARGIN,
                top: Inches = FIG_TOP, max_w: Inches = Inches(8.2)):
    w, h = _fit_figure(panel, target_h, max_w)
    slide.shapes.add_picture(str(panel), left, top, height=h, width=w)
    return w, h


def _add_caption(slide, intro, left: Inches, top: Inches, width: Inches,
                 height: Inches, body_size=Pt(14), caveat: str | None = None) -> None:
    cap = slide.shapes.add_textbox(left, top, width, height)
    ctf = cap.text_frame
    ctf.word_wrap = True
    for i, para_text in enumerate(intro):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.space_after = Pt(9)
        p.alignment = PP_ALIGN.LEFT
        _add_run_with_italics(p, para_text, body_size, INK)
    if caveat:
        p = ctf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = caveat
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = MUTED


def _add_table(slide, headers, rows, left, top, width, height, col_fracs=None,
               font=Pt(11)) -> None:
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gfx.table
    if col_fracs:
        for idx, frac in enumerate(col_fracs):
            table.columns[idx].width = Emu(int(width * frac))
    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_FILL
        for p in cell.text_frame.paragraphs:
            p.font.size = font
            p.font.bold = True
            p.font.color.rgb = INK
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = font
                p.font.color.rgb = INK


def _blank_slide(prs: Presentation):
    # layout 6 ("Blank") in python-pptx's default template has no placeholders
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# derived TCF feature table (read from the CSV built in the audit phase)
# ---------------------------------------------------------------------------
def _tcf_feature_rows():
    """Collapse roi_feature_summary_hp3_tcf.csv to one row per ROI position."""
    by_pos: dict[int, dict] = {}
    with open(TCF_TABLE_CSV, newline="") as fh:
        for row in csv.DictReader(fh):
            pos = int(row["position"])
            d = by_pos.setdefault(pos, {"sub": row["substitution"], "ranks": {}})
            d["ranks"][int(row["rank"])] = row
    out = []
    for pos in sorted(by_pos):
        d = by_pos[pos]
        r = d["ranks"]
        def cell(k):
            x = r[k]
            return f"{x['feature_id']} ({x['category']})\n~{float(x['wt_activation']):.2f} / {float(x['evolved_activation']):.2f}"
        out.append([d["sub"], cell(1), cell(2), cell(3)])
    return out


# ---------------------------------------------------------------------------
# case-study slides
# ---------------------------------------------------------------------------
def slide_overview(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "HP3 \u2192 HP3.1 overview", "BACKGROUND", MUTED)
    _add_subtitle(slide, "Two structural proteins, two kinds of change")
    # annotated T4 schematic on the left, showing where TCF and gp8 sit
    fig_w, _ = _add_figure(slide, T4_OVERVIEW, Inches(4.7), max_w=Inches(4.6))
    cap_left = MARGIN + fig_w + Inches(0.3)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "*Escherichia* phage HP3 (NC_041919.1) and its evolved relative HP3.1 "
        "(OK275722.1) differ in two tail-apparatus proteins we track here.",
        "*Tail collar fiber (TCF)* \u2014 516-aa homotrimer (HP3 YP_010228801.1 / "
        "HP3.1 URQ01387.1). Three point substitutions, no indels: *Y444H, K464R, "
        "Y465H*.",
        "*Baseplate wedge gp8* \u2014 a T4-like gp8 (HP3 YP_010228805.1, 334 aa / "
        "HP3.1 URQ01383.1, 337 aa). Not a substitution: a *3-residue DPN "
        "tandem-repeat insertion* at evolved positions 290\u2013292.",
        "Both are read residue-by-residue through ESMC-6B and its layer-60 SAE, "
        "then compared WT vs evolved in SAE-feature space and on AlphaFold-3 "
        "structures. The contrast built over the next slides: the TCF substitutions "
        "sit at a conserved interface *without* changing the SAE feature profile, "
        "while the gp8 insertion *relocates* an interaction feature.",
    ]
    _add_caption(slide, intro, cap_left, FIG_TOP, cap_w,
                 Inches(3.3), body_size=Pt(14))
    headers = ["Protein", "Biological unit", "HP3 \u2192 HP3.1 change", "Length"]
    rows = [
        ["Tail collar fiber", "Homotrimer (A/B/C)", "Y444H, K464R, Y465H (3 substitutions)", "516 \u2192 516 aa"],
        ["Baseplate wedge gp8", "T4-like wedge subunit", "insDPN at 290\u2013292 (tandem duplication)", "334 \u2192 337 aa"],
    ]
    _add_table(slide, headers, rows, cap_left, Inches(5.25), cap_w,
               Inches(1.2), col_fracs=[0.22, 0.22, 0.40, 0.16], font=Pt(11))
    _add_footer(slide, "Phage schematic: PhageExterior.svg, Wikimedia Commons "
                       "(Adenosine, CC BY-SA 3.0), circles added.", 4)
    slide.notes_slide.notes_text_frame.text = (
        "Provenance. TCF: both 516 aa, no indels; three substitutions confirmed by direct "
        "sequence diff and matching data/hp3_to_hp3e.txt (Y444H, K464R, Y465H). gp8: WT 1-289 "
        "== evolved 1-289; evolved gains DPN at 290-292; WT 290-334 == evolved 293-337 "
        "(data/gp8_alignment.csv, direct string alignment) \u2014 the insertion sits between WT "
        "residue 289 (N) and WT 290 (V). AF3 monomers: CA count matches FASTA length exactly "
        "in every model (gp8 334/337, TCF 516/516); numbering is contiguous 1..N so sequence "
        "position == auth_seq_id. Model: ESMC-6B, layer-60 SAE, 16,384 features, unsupervised "
        "(same method as the two opener slides)."
    )


def slide_tcf_structure(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "TCF trimer structure", "TAIL COLLAR FIBER", TCF_ACCENT)
    _add_subtitle(slide, "Three substitutions, all at inter-chain contacts",
                  color=TCF_ACCENT)
    fig_w, _ = _add_figure(slide, CS / "tcf_3mer_roi_clusters_legend.png",
                           Inches(5.1), max_w=Inches(6.7))
    cap_left = MARGIN + fig_w + Inches(0.3)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "The tail collar fiber's biological unit is a homotrimer \u2014 an elongated "
        "fiber (surface view), SAE-cluster colored (legend below the structure).",
        "All three substituted positions (*Y444H, K464R, Y465H*) sit at *direct "
        "inter-chain contact* in both variants: measured CA/side-chain minimum "
        "distances of *2.85\u20133.28 \u00c5* to a neighboring chain.",
        "So these are not surface-exposed passenger mutations \u2014 they are at the "
        "trimer interface, where a change can tune subunit packing.",
    ]
    _add_caption(slide, intro, cap_left, FIG_TOP, cap_w, Inches(3.0), body_size=Pt(14))
    headers = ["Position", "WT \u2192 Evolved", "Contact (evolved)"]
    rows = [
        ["444", "Y \u2192 H", "2.87\u20132.97 \u00c5"],
        ["464", "K \u2192 R", "2.85\u20133.11 \u00c5"],
        ["465", "Y \u2192 H", "2.85\u20133.11 \u00c5"],
    ]
    _add_table(slide, headers, rows, cap_left, Inches(4.9), cap_w, Inches(1.4),
               col_fracs=[0.24, 0.40, 0.36], font=Pt(12))
    _add_footer(slide, None, 5)
    slide.notes_slide.notes_text_frame.text = (
        "Interface contacts from src/tcf_interface_contacts.py on the AF3 homotrimers "
        "(analysis/structures/tcf_interface_contacts.csv). Per-chain nearest-neighbor min "
        "distances: WT 444 3.04-3.07 A, 464 3.01-3.04 A, 465 3.21-3.28 A; evolved 444 "
        "2.87-2.97 A, 464 2.85-3.11 A, 465 2.85-3.11 A \u2014 all three positions are at direct "
        "inter-chain contact in both variants. The evolved 3mer was predicted AFTER the "
        "2026-07-21 Y444H correction (all three chains confirmed H444/R464/H465); the older "
        "evolved *monomer* still carries Y at 444 and should not be used for 444 side-chain "
        "geometry. Surface coloring is by SAE residue cluster, not by ROI."
    )


def slide_tcf_features(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "TCF SAE features", "TAIL COLLAR FIBER", TCF_ACCENT)
    _add_subtitle(slide, "The substitutions leave the feature profile unchanged",
                  color=TCF_ACCENT)
    fig_w, _ = _add_figure(slide, FIG / "roi_heatmap_hp3_tcf_sae.png",
                           Inches(4.7), max_w=Inches(5.4))
    cap_left = MARGIN + fig_w + Inches(0.3)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "Pairwise cosine distance between ROI residues in SAE space. Each WT/evolved "
        "pair (Y444/H444, K464/R464, Y465/H465) sits at *~0.00* \u2014 the substitution "
        "barely moves the residue in feature space.",
        "The top three SAE features are *identical* across WT and evolved at all "
        "three positions (feature 10240 disorder, 13260 phage tail-fiber motif, "
        "15023 structural motif), with near-identical activations.",
        "Reading: ESMC/SAE sees these three positions as the *same* conserved "
        "tail-fiber loop before and after \u2014 the sequence changes conservatively "
        "enough that the learned feature description is preserved.",
    ]
    _add_caption(slide, intro, cap_left, FIG_TOP, cap_w, Inches(2.6), body_size=Pt(13))
    headers = ["Substitution", "Rank 1 (Disorder)", "Rank 2 (Comp. bias)", "Rank 3 (Struct.)"]
    rows = _tcf_feature_rows()
    _add_table(slide, headers, rows, cap_left, Inches(4.55), cap_w, Inches(2.1),
               col_fracs=[0.16, 0.30, 0.30, 0.24], font=Pt(10))
    _add_footer(slide, None, 6)
    slide.notes_slide.notes_text_frame.text = (
        "Table cells: feature id (SAE category), then WT / evolved activation. Source "
        "analysis/features/roi_feature_summary_hp3_tcf.csv (derived from "
        "hp3_whole_protein_features_top3_tcf.csv). Rank1 = feature 10240 'Polar loop hotspots; "
        "S/T-rich disorder' (~1.25-1.34); rank2 = 13260 'Phage tail fiber distal motifs' "
        "(~0.70-0.73); rank3 = 15023 'Active-site loop/turn scaffolds' (~0.22-0.64). Feature id "
        "and rank are shared WT<->evolved at every position. This is the direct contrast with "
        "gp8, where the insertion relocates the rank-3 interaction feature. Off-diagonal "
        "structure (444 ~0.08-0.09 from 464/465; 464<->465 ~0.01-0.02) reflects real "
        "position-to-position differences, independent of the WT/evolved swap. ESMC is "
        "sequence-only; these are learned sequence-context labels, not measured 3D contacts."
    )


def slide_tcf_head_axis(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "TCF head domain", "TAIL COLLAR FIBER", TCF_ACCENT)
    _add_subtitle(slide, "The three ROI positions cluster at the fiber tip",
                  color=TCF_ACCENT)
    # three-view panel of the evolved (HP3e) head domain; ROI in pink
    n = len(TCF_HEAD_VIEWS)
    gap = Inches(0.3)
    col_w = (SLIDE_W.inches - 2 * MARGIN.inches - (n - 1) * gap.inches) / n  # column pitch
    panel_max_h = 3.35   # cap height so the caption band clears the footer
    panel_top = Inches(1.85)
    lab_h = Inches(0.30)
    for i, (path, label) in enumerate(TCF_HEAD_VIEWS):
        left = MARGIN.inches + i * (col_w + gap.inches)
        side = min(col_w, panel_max_h)   # square images
        w = h = Inches(side)
        x = Inches(left + (col_w - side) / 2)   # center in the column
        slide.shapes.add_picture(str(path), x, panel_top, width=w, height=h)
        lb = slide.shapes.add_textbox(Inches(left), Inches(panel_top.inches + side + 0.05),
                                      Inches(col_w), lab_h)
        p = lb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = TCF_ACCENT
    # single caption spanning the width, below the panels
    cap_top = Inches(panel_top.inches + min(col_w, panel_max_h) + lab_h.inches + 0.14)
    cap = slide.shapes.add_textbox(MARGIN, cap_top, SLIDE_W - 2 * MARGIN, Inches(1.0))
    ctf = cap.text_frame; ctf.word_wrap = True
    _add_run_with_italics(
        ctf.paragraphs[0],
        "Head domain of the evolved (HP3.1) TCF homotrimer, surface view; the three ROI "
        "positions (*H444 / R464 / H465*, pink) from three orientations. Down the 3-fold "
        "axis (left) they converge at the trimer center where the subunits meet; the "
        "tilted and side views place that patch at the distal tip \u2014 the inter-chain "
        "contact, matching the unchanged SAE feature profile.",
        Pt(14), INK)
    _add_footer(slide, "Evolved (HP3.1) TCF head; ROI H444/R464/H465 in pink.", 7)
    slide.notes_slide.notes_text_frame.text = (
        "Surface views of the evolved (HP3.1 / HP3e) TCF homotrimer head domain with the ROI "
        "residues 444/464/465 (H444/R464/H465 in the evolved variant) highlighted in pink, "
        "shown in three orientations: straight down the 3-fold axis from the distal tip "
        "(full-on), tilted, and side-on. The full-on view shows all three ROI copies (one per "
        "chain) meeting at the center of the trimer where the three subunits pack together; the "
        "side/tilted views show the same patch sitting at the distal tip of the fiber. This "
        "slide REPLACES the earlier down-axis SAE-cluster surface (and, before that, the "
        "three-residue interface close-up) per the presenter's request; the three panels are "
        "the presenter-supplied renders tcf_head_hp3e_{full_on,tilted,side}.png. Takeaway for "
        "the audience: the three substitutions are all at the inter-chain contact at the fiber "
        "tip, consistent with the SAE result that their feature description is unchanged."
    )


def slide_gp8_dpn(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "gp8 SAE features", "BASEPLATE gp8", GP8_ACCENT)
    _add_subtitle(slide, "A DPN insertion relocates the interaction site",
                  color=GP8_ACCENT)
    fig_w, _ = _add_figure(slide, CS / "gp8_2mer_evolved_roi_clusters.png",
                           Inches(4.5), max_w=Inches(5.0))
    cap_left = MARGIN + fig_w + Inches(0.3)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "gp8 gains a second *DPN* copy: WT D287-P288-N289 is followed by V290; the "
        "evolved chain inserts a fresh D290-P291-N292 before the original V "
        "(now V293).",
        "Ranks 1\u20132 are identical across 287\u2013293 in both variants (disorder + a "
        "weak structural module) \u2014 the region reads as a flexible surface loop.",
        "*Rank 3 is where they diverge, at the same sequence position:* the first "
        "DPN copy always reads \u201cInteraction site\u201d (feature 7819); V290 (WT) and "
        "the inserted second copy (evolved) both read \u201cDomain\u201d (feature 3038).",
        "*Conclusion: the duplication relocates the interaction site rather than "
        "duplicating it* \u2014 ESMC/SAE only ever calls one DPN copy "
        "interaction-relevant, and it is always the first.",
    ]
    _add_caption(slide, intro, cap_left, FIG_TOP, cap_w, Inches(2.9), body_size=Pt(12.5),
                 caveat="ESMC is sequence-only \u2014 learned sequence-context labels, not "
                        "measured 3D contacts (see next two slides).")
    headers = ["Region", "WT", "Evolved", "Rank 3"]
    rows = [
        ["First DPN copy (invariant)", "D287-P288-N289", "D287-P288-N289", "7819 \u2014 Interaction site"],
        ["V290 / relocated 2nd copy", "V290", "D290-P291-N292-V293", "3038 \u2014 Domain"],
    ]
    _add_table(slide, headers, rows, cap_left, Inches(5.25), cap_w, Inches(1.2),
               col_fracs=[0.28, 0.20, 0.28, 0.24], font=Pt(10))
    _add_footer(slide, None, 8)
    slide.notes_slide.notes_text_frame.text = (
        "Evolved 290-293 rank 4 = feature 7819 (Interaction site) on all four residues "
        "(~0.34-0.37) \u2014 a secondary echo, not a second strong site. Rank 5 varies per residue: "
        "D290->7070 'Phage tail attachment segments' (0.295); P291->3509 'C-terminal "
        "helix-to-IDR transition' (0.303); N292->12804 'Cysteine-rich LU/CRD domains' (0.303); "
        "V293->3509 (0.316). WT V290 rank4/5 = 7819 (0.359) / 7070 (0.309). Insertion coords "
        "confirmed by data/gp8_alignment.csv: DPN at evolved 290-292, between WT 289 (N) and WT "
        "290 (V). Ranks 1-2 across 287-293: feature 15659 disorder (~0.60-0.65) + 15513 "
        "structural motif (~0.40-0.44)."
    )


def slide_gp8_wedge(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "gp8 in the T4 wedge", "BASEPLATE gp8", GP8_ACCENT)
    _add_subtitle(slide, "Its real baseplate-wedge context (PDB 9MKB)",
                  color=GP8_ACCENT)
    fig_w, _ = _add_figure(slide, CS / "gp8_evolved_in_wedge_context.png",
                           Inches(5.0), max_w=Inches(5.4))
    cap_left = MARGIN + fig_w + Inches(0.3)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "The evolved gp8 (SAE-cluster colored) placed inside its authentic "
        "neighborhood: the cryo-EM T4 baseplate wedge (PDB *9MKB*), not an isolated "
        "monomer.",
        "Confirms gp8 is a genuine T4 gp8 wedge homolog \u2014 CA-centroid overlay onto "
        "9MKB chain ZE differs by only *0.30 \u00c5*.",
        "Wedge stoichiometry recovered exactly as T4 biology predicts: *2\u00d7 gp8 + "
        "gp7 + gp6* (background context, de-emphasized).",
        "The ROI (287\u2013293, black spheres) sits at the boundary facing the second "
        "gp8 copy \u2014 near a real inter-subunit interface, not buried mid-fold.",
        "Sets up the open question: does the *inserted* second DPN copy make its own "
        "contact, or is the loop simply repositioned?",
    ]
    _add_caption(slide, intro, cap_left, FIG_TOP, cap_w, FIG_BOTTOM_MAX - FIG_TOP,
                 body_size=Pt(13.5))
    _add_footer(slide, "Overlay onto cryo-EM T4 baseplate, PDB 9MKB.", 9)
    slide.notes_slide.notes_text_frame.text = (
        "9MKB is the cryo-EM T4 portal-neck-tail complex (541 chains). A CA-CA contact search "
        "(<10 A) from chain ZE (which overlays our AF3 gp8 almost exactly) against all other "
        "chains finds exactly 6 true neighbors: ZD (2nd gp8), YC/Yc (gp7), eB/eU/f9 (gp6) \u2014 "
        "matching the known wedge stoichiometry (2x gp8 + gp7 + gp6 per wedge). The other 534 "
        "chains (tail tube/sheath, capsid, portal, fibers) were dropped, not just hidden. "
        "Built by src/gp8_neighbor_context_hp3.py; source overlay "
        "analysis/structures/gp8_overlay/9mkb_gp8_neighbors.cif."
    )


def slide_gp8_overlay(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "gp8 on the T4 tail", "BASEPLATE gp8", GP8_ACCENT)
    _add_subtitle(slide, "Does the duplication move the contact point?",
                  color=GP8_ACCENT)
    fig_w, _ = _add_figure(slide, CS / "t4_tail_overlay_montage_3x3.png",
                           Inches(5.1), max_w=Inches(6.0))
    cap_left = MARGIN + fig_w + Inches(0.3)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "3\u00d73 comparison. Columns: *T4 gp8* (reference), *HP3* (WT), *HP3e* "
        "(evolved). Rows: axial, transverse, cross-section views of gp8 overlaid on "
        "the real T4 tail.",
        "Across all three angles, the DPN duplication visibly *shifts the "
        "first-contact loop* relative to WT \u2014 independent structural confirmation of "
        "the SAE finding: the footprint does not get bigger, it moves.",
        "Why duplicate the whole *DPN* triplet, not a partial insertion? DPN is a "
        "proline-anchored turn unit \u2014 duplicating the full triplet preserves the "
        "local turn geometry while shifting the downstream position.",
        "*Working conclusion: the insertion repositions the binding location*, "
        "consistent with the tropism shift seen in HP3e.",
    ]
    _add_caption(slide, intro, cap_left, FIG_TOP, cap_w, FIG_BOTTOM_MAX - FIG_TOP,
                 body_size=Pt(13.5),
                 caveat="Still open: whether the inserted copy independently contacts the "
                        "tail (vs. riding along) is not resolved by the sequence-only SAE or "
                        "this overlay \u2014 flagged for a direct contact measurement.")
    _add_footer(slide, "AF3 gp8 models overlaid on the T4 tail (9MKB); 3 views \u00d7 "
                       "{T4, HP3, HP3e}.", 10)
    slide.notes_slide.notes_text_frame.text = (
        "Panels built from AF3 models of HP3 (WT) and HP3e (evolved) gp8, structurally aligned "
        "onto the real T4 tail (9MKB) coordinate frame, then rendered from three fixed camera "
        "angles per variant so the same viewpoint is directly comparable across T4/HP3/HP3e. "
        "Generator src/cluster_structure_hp3.py; individual-view PNGs also in "
        "analysis/figures/cluster_structures_hp3/ (t4_tail_overlay_{axial,transverse,xsection}_"
        "{t4,hp3,hp3e}.png)."
    )


def slide_synthesis(prs):
    slide = _blank_slide(prs)
    _add_header_bars(slide, "Synthesis", "CONCLUSIONS", MUTED)
    _add_subtitle(slide, "Two proteins, two mechanisms \u2014 one readable method")
    intro = [
        "*Tail collar fiber (substitution at a conserved interface).* Y444H / K464R "
        "/ Y465H all sit at direct inter-chain contact (2.85\u20133.28 \u00c5), yet the SAE "
        "feature profile is unchanged (WT\u2013evolved cosine ~0.00; identical top-3 "
        "features). The changes are conservative tuning at a preserved site.",
        "*Baseplate gp8 (insertion that relocates).* The DPN tandem duplication does "
        "not add a second interaction site \u2014 SAE calls only the first DPN copy "
        "interaction-relevant, and the T4-tail overlay shows the contact loop moving "
        "rather than growing.",
        "*Why this matters for the method.* In both cases the call is traceable to "
        "named SAE features and their activations \u2014 the audit the labels-come-in-"
        "post-hoc workflow was built to give: *these named features drove this call.*",
        "*Open, testable questions.* (1) Does the inserted gp8 DPN copy make an "
        "independent inter-chain contact? Needs a direct distance measurement on the "
        "aligned overlay. (2) Do the TCF interface substitutions measurably change "
        "trimer packing or host-binding affinity?",
    ]
    _add_caption(slide, intro, MARGIN, FIG_TOP, SLIDE_W - 2 * MARGIN,
                 FIG_BOTTOM_MAX - FIG_TOP, body_size=Pt(16))
    _add_footer(slide, None, 11)
    slide.notes_slide.notes_text_frame.text = (
        "Closing contrast for the audience: the SAE didn't just label things, it distinguished "
        "two evolutionary moves that look superficially similar (both at the tail apparatus, "
        "both near interfaces) but are mechanistically different \u2014 a conservative substitution "
        "that preserves the feature description vs. an insertion that relocates a feature. Both "
        "conclusions are backed by AF3 structure (interface contacts; wedge/tail overlays) as "
        "well as the sequence-only SAE, so the two lines of evidence agree. Deferred audit point "
        "from the opener lands here: traceability to named features == the audit."
    )


def _feature_grid(slide, left: Inches, top: Inches, rows: int, cols: int,
                   cell: Inches, gap: Inches, lit: set, dim_color: RGBColor,
                   lit_color: RGBColor) -> None:
    """Draw a rows x cols grid of small squares; cells in `lit` (set of (r, c))
    render in `lit_color`, everything else in `dim_color` \u2014 the "sparse
    features" motif shared by both opener slides."""
    for r in range(rows):
        for c in range(cols):
            x = Inches(left.inches + c * (cell.inches + gap.inches))
            y = Inches(top.inches + r * (cell.inches + gap.inches))
            color = lit_color if (r, c) in lit else dim_color
            _rect(slide, x, y, cell, cell, color)


def slide_opener_black_box(prs):
    """Opener 1/2: SAEs as a validated interpretability method (dark ground)."""
    slide = _blank_slide(prs)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK_BG)
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.4), SLIDE_W - 2 * MARGIN, Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Opening the black box"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    _small_caps(r)
    _add_subtitle(slide, "Sparse autoencoders (SAEs), proven on frontier LLMs", color=WHITE)

    # left: one dense block standing in for entangled, unlabeled activations
    block_left, block_top = HX(0.6), Inches(1.85)
    block_w, block_h = HX(3.2), Inches(2.9)
    _rect(slide, block_left, block_top, block_w, block_h, DIM_CELL)
    lb = slide.shapes.add_textbox(block_left, Inches(block_top.inches + block_h.inches + 0.1),
                                  block_w, Inches(0.5))
    p = lb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "activations (entangled)"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = WHITE

    # right: same footprint, mostly dark, a few lit \u2014 sparsity after the SAE
    grid_left = HX(7.4)
    grid_top = Inches(1.85)
    cell, gap = Inches(0.36), Inches(0.08)
    _feature_grid(slide, grid_left, grid_top, rows=6, cols=6, cell=cell, gap=gap,
                  lit={(1, 1), (2, 4), (4, 2), (5, 5)},
                  dim_color=DIM_CELL, lit_color=LIT_CELL)
    grid_w = Inches(6 * cell.inches + 5 * gap.inches)
    lb = slide.shapes.add_textbox(grid_left, Inches(grid_top.inches + 6 * (cell.inches + gap.inches) + 0.05),
                                  grid_w, Inches(0.5))
    p = lb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "SAE features (sparse, most off)"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = WHITE

    intro = [
        "*The problem.* Large models pack far more concepts than they have neurons, so "
        "no single neuron is individually readable \u2014 the model works, but it can't be "
        "inspected.",
        "*The fix.* A sparse autoencoder is a wide, simple layer bolted onto a *frozen* "
        "model. It re-expresses the model's internal activations as a large dictionary "
        "of features, almost all switched off at any moment, each one human-interpretable.",
        "*Proof points.* Anthropic decomposed Claude into millions of interpretable "
        "features and could steer behavior by amplifying or suppressing them; OpenAI did "
        "the same at roughly 16 million features on GPT-4 \u2014 a mature, reproducible method.",
    ]
    cap = slide.shapes.add_textbox(MARGIN, Inches(5.45), SLIDE_W - 2 * MARGIN, Inches(1.5))
    ctf = cap.text_frame
    ctf.word_wrap = True
    for i, para_text in enumerate(intro):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.space_after = Pt(4)
        _add_run_with_italics(p, para_text, Pt(13), WHITE)
    _add_footer(slide, None, 2)


def slide_opener_esmc(prs):
    """Opener 2/2: the same method applied to ESMC; labels come in post-hoc."""
    slide = _blank_slide(prs)
    _add_header_bars(slide, "Same method, our proteins", "METHOD", MUTED)
    _add_subtitle(slide, "ESMC + the identical sparse-autoencoder recipe")

    grid_left, grid_top = HX(0.6), Inches(1.95)
    cell, gap = Inches(0.42), Inches(0.1)
    lit = {(1, 1), (2, 4), (4, 2), (5, 5)}
    _feature_grid(slide, grid_left, grid_top, rows=6, cols=6, cell=cell, gap=gap,
                  lit=lit, dim_color=LIGHT_CELL, lit_color=LIT_CELL)
    grid_w = Inches(6 * cell.inches + 5 * gap.inches)
    labels = ["binding site", "membrane assoc.", "catalytic motif", "structural motif"]
    for (r, c), label in zip(sorted(lit), labels):
        x = Inches(grid_left.inches + c * (cell.inches + gap.inches) - 0.5)
        y = Inches(grid_top.inches + r * (cell.inches + gap.inches) + cell.inches + 0.03)
        lb = slide.shapes.add_textbox(x, y, Inches(1.4), Inches(0.3))
        p = lb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.color.rgb = INK

    cap_left = Inches(grid_left.inches + grid_w.inches + 0.5)
    cap_w = SLIDE_W - cap_left - MARGIN
    intro = [
        "The identical technique applied to ESMC (a protein language model) yields a "
        "labeled catalog of biology instead of language: *16,384 features*, learned "
        "unsupervised, mapping onto pre-existing biology \u2014 binding sites, membrane "
        "association, catalytic function, structural motifs.",
        "*The training objective:* L = \u2016a \u2212 Dz\u2016\u00b2 + \u03bb\u2016z\u2016\u2081,  "
        "z = ReLU(Wa + b) \u2014 a is the frozen ESM activation for one residue, Dz is its "
        "reconstruction from a sparse feature dictionary D, and the \u2016\u00b7\u2016\u2081 "
        "penalty forces most features off, which is what makes each one readable.",
        "*The dictionary D is learned with no labels* \u2014 the biological descriptions are "
        "attached only afterward, matched against known function. Labels come in "
        "post-hoc; the case study that follows shows what that traceability buys.",
    ]
    _add_caption(slide, intro, cap_left, Inches(1.95), cap_w, Inches(4.6), body_size=Pt(13.5))
    _add_footer(slide, None, 3)


def main() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_opener_black_box(prs)
    slide_opener_esmc(prs)
    slide_overview(prs)
    slide_tcf_structure(prs)
    slide_tcf_features(prs)
    slide_tcf_head_axis(prs)
    slide_gp8_dpn(prs)
    slide_gp8_wedge(prs)
    slide_gp8_overlay(prs)
    slide_synthesis(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(ROOT)}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
